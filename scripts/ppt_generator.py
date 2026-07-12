#!/usr/bin/env python3
"""
环保执法 PPT 生成器
基于 python-pptx 生成培训课件和执法汇报 PPT
"""

import os, json
from datetime import date
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

GREEN = RGBColor(0x1B, 0x43, 0x32)
LIGHT_GREEN = RGBColor(0xD8, 0xE8, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)

prs = None

def new_prs():
    global prs
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs

def add_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GREEN
    bg.line.fill.background()

    # 装饰线
    line = slide.shapes.add_shape(1, Cm(2), Cm(7), Cm(29.867), Cm(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    tx = slide.shapes.add_textbox(Cm(2), Cm(7.5), Cm(29.867), Cm(4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        tx2 = slide.shapes.add_textbox(Cm(2), Cm(12), Cm(29.867), Cm(2))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = LIGHT_GREEN
        p2.alignment = PP_ALIGN.CENTER

    tx3 = slide.shapes.add_textbox(Cm(2), Cm(16), Cm(29.867), Cm(2))
    tf3 = tx3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f'济南市生态环境局 | {date.today().strftime("%Y年%m月")}'
    p3.font.size = Pt(14)
    p3.font.color.rgb = LIGHT_GREEN
    p3.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GREEN
    bg.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, Cm(8), Cm(5), Cm(3.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_GREEN
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Cm(6), Cm(8), Cm(25), Cm(4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide

def add_content_slide(title, bullets, note=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 顶部条
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Cm(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    # 标题
    tx = slide.shapes.add_textbox(Cm(1), Cm(0.4), Cm(31), Cm(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 内容
    content = slide.shapes.add_textbox(Cm(1.5), Cm(3.2), Cm(30.5), Cm(13))
    tf = content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'• {bullet}'
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(6)
        p.space_after = Pt(4)

    if note:
        tx2 = slide.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(30), Cm(1.5))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'📌 {note}'
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.font.italic = True

    # 页码
    pg = slide.shapes.add_textbox(Cm(30), Cm(17.5), Cm(3), Cm(1))
    tf = pg.text_frame
    p = tf.paragraphs[0]
    p.text = f'{prs.slides.index(slide) + 1} / {len(prs.slides)}'
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT
    return slide

def add_table_slide(title, headers, rows, note=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Cm(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Cm(1), Cm(0.4), Cm(31), Cm(2))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(n_rows, n_cols, Cm(1.5), Cm(3.2), Cm(30.5), Cm(12)).table

    # 设置列宽
    col_w = Cm(30.5 / n_cols)
    for i in range(n_cols):
        for row in table.rows:
            row.cells[i].width = col_w

    # 标题行
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # 数据行
    for r_i, row_data in enumerate(rows, 1):
        for c_i, text in enumerate(row_data):
            cell = table.cell(r_i, c_i)
            cell.text = str(text)
            bg = LIGHT_GREEN if r_i % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK

    if note:
        tx2 = slide.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(30), Cm(1.5))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = f'📌 {note}'
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.font.italic = True
    return slide

def generate_training_ppt(topic='行政处罚程序', output_path=None):
    """生成环保执法培训 PPT"""
    new_prs()
    add_title_slide(f'环保执法实务培训\n{topic}', '济南市生态环境局执法培训课程')

    topics_map = {
        '行政处罚程序': _build_admin_penalty_slides,
        '典型案例分析': _build_case_study_slides,
        '法律法规汇编': _build_law_compilation_slides,
        '执法风险防范': _build_risk_prevention_slides,
        '案件办理流程': _build_case_process_slides,
    }
    builder = topics_map.get(topic, _build_admin_penalty_slides)
    builder()

    if not output_path:
        output_path = os.path.join(os.path.dirname(__file__), '..', 'outputs', f'培训_{topic}_{date.today().strftime("%Y%m%d")}.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path

def _build_admin_penalty_slides():
    add_section_slide('一、行政处罚程序概述')
    add_content_slide('行政处罚程序总览', [
        '简易程序：事实清楚、证据确凿，执法人员当场作出处罚决定',
        '一般程序：立案→调查→告知→决定→送达→执行（标准流程）',
        '听证程序：当事人要求听证的，应在3日内提出，7日前通知听证时间地点',
        '执行程序：当事人自觉履行 → 逾期加处罚款 → 申请法院强制执行',
        '结案归档：处罚决定执行完毕后5日内填写结案报告，归档保存',
    ], '注：一般程序是最常用的行政处罚程序，约占80%的案件')

    add_section_slide('二、立案条件与管辖')
    add_content_slide('立案条件', [
        '① 发现或接到违法行为线索',
        '② 有明确违法嫌疑人',
        '③ 属于本机关管辖范围',
        '④ 依法应当给予行政处罚',
        '⑤ 违法行为仍在追诉时效内（一般为2年）',
    ], '管辖原则：违法行为发生地 > 结果发生地 > 嫌疑人住所地')
    add_content_slide('不予立案情形', [
        '① 违法行为轻微并及时纠正，没有造成危害后果的',
        '② 超过行政处罚时效（违法行为终了之日起2年）',
        '③ 不属于本机关管辖范围',
        '④ 涉嫌犯罪需移送司法机关的（不适用行政处罚）',
        '⑤ 其他依法不予处罚的情形',
    ])

    add_section_slide('三、调查取证规范')
    add_content_slide('现场检查要点', [
        '亮证执法：出示行政执法证件，不少于2名执法人员',
        '现场笔录：客观记录时间、地点、人物、事实、证据',
        '视听资料：现场拍照、录像，记录原始状态',
        '采样监测：依法采样，填写采样记录，委托有资质机构检测',
        '书证调取：环评批复、排污许可证、运行台账、用电记录等',
    ], '证据"三性"要求：真实性、合法性、关联性')
    add_content_slide('询问调查技巧', [
        '告知程序：告知被询问人权利义务，如实回答义务',
        '单刀直入：直接针对违法行为核心问题提问',
        '细节追问：对时间、地点、方式、数量等关键细节深挖',
        '证据印证：通过书面证据与口供相互印证',
        '禁止事项：不得威胁、引诱、欺骗被询问人',
    ])

    add_section_slide('四、行政处罚裁量')
    add_content_slide('裁量基本原则', [
        '过罚相当：处罚种类和幅度与违法行为的性质、情节相适应',
        '公平公正：同类案件同等处理，不因当事人身份区别对待',
        '程序正当：严格遵循法定程序，保障当事人陈述申辩权',
        '教育与处罚相结合：处罚不是目的，促进守法才是目标',
    ], '山东省生态环境厅已发布裁量基准表，对主要违法行为明确了细化标准')
    add_table_slide('裁量因素参考表', ['情形类别', '从轻因素', '从重因素'], [
        ['主观过错', '主动停止/消除危害', '故意违法/恶意逃避监管'],
        ['违法持续时间', '< 30天', '> 180天'],
        ['配合调查', '如实陈述/主动提交证据', '拒绝调查/销毁证据'],
        ['整改态度', '积极整改/消除污染', '拒不整改/继续违法'],
        ['危害后果', '未造成明显危害', '造成人群健康损害/重大环境污染'],
    ])

    add_section_slide('五、典型案例警示')
    add_table_slide('违法类型与处罚对照', ['违法类型', '主要法条', '罚款区间', '典型后果'], [
        ['私设暗管偷排', '水法第83条', '10-100万', '停产停业+强制执行'],
        ['超标排放大气', '大气法第99条', '10-100万', '责令限制生产'],
        ['危废非法倾倒', '固废法第112条', '10-100万', '移送公安机关（危废≥3吨入刑）'],
        ['未批先建', '环评法第31条', '总投资额1-5%', '责令停建+罚款'],
        ['无证排污', '排污许可条例33条', '20-100万', '责令限制生产/停业'],
        ['监测数据造假', '大气法第100条', '10-100万', '移送公安（破坏计算机信息系统罪）'],
    ], '提示：造成严重污染环境的，依据刑法第338条可处3-7年有期徒刑')
    add_content_slide('执法风险提示', [
        '⚠ 程序风险：超期办案（一般程序应在90日内办结，可延长至180日）',
        '⚠ 证据风险：取证实不规范导致证据无效（如采样不规范、无见证人）',
        '⚠ 裁量风险：畸轻畸重面临行政复议或诉讼被撤销风险',
        '⚠ 执行风险：超期申请法院强制执行（2年内）',
        '✓ 防范措施：严格程序、规范取证、集体讨论、充分告知',
    ])

def _build_case_study_slides():
    add_section_slide('典型案例深度剖析')
    add_content_slide('案例一：某化工公司私设暗管偷排案', [
        '案情：执法人员深夜突击检查，发现厂区雨水管道接入市政管网，经检测COD超标18倍',
        '定性：私设暗管逃避监管方式排放水污染物',
        '处罚：罚款52万元 + 责令停产整治',
        '要点：暗管不仅限于地下管道，雨水管道排污水同样构成私设暗管',
        '延伸：当事人被以污染环境罪移送，判处有期徒刑1年6个月',
    ], '关键证据：夜间突击检查 + 暗管接入点确认 + 污水COD检测报告')
    add_content_slide('案例二：在线监测数据造假案', [
        '案情：第三方运维人员对COD在线监测设备参数篡改，导致上传数据严重失真',
        '定性：篡改监测数据，以逃避监管方式排放污染物',
        '处罚：对排污单位罚款45万元；对运维公司罚款20万元',
        '要点：排污单位和第三方运维公司可分别处罚，构成共同违法',
        '延伸：运维人员以破坏计算机信息系统罪判处有期徒刑2年',
    ], '监测数据造假入刑门槛：无门槛（后果严重即可），不必等超标数据积累')

def _build_law_compilation_slides():
    add_section_slide('常用法律法规汇编')
    add_table_slide('国家级核心法规一览', ['法规名称', '施行日期', '核心罚则'], [
        ['环境保护法', '2015.01.01', '第63条：停产停业/移送拘留；第69条：公益诉讼'],
        ['大气污染防治法', '2016.01.01', '第99/100条：10-100万罚款（超标/数据造假）'],
        ['水污染防治法', '2018.01.01', '第83条：私设暗管/超标排放，10-100万罚款'],
        ['固体废物污染环境防治法', '2020.09.01', '第112条：危废非法处置，第123条：按日计罚'],
        ['环境影响评价法', '2016.09.01', '第31条：未批先建，总投资额1-5%罚款'],
        ['排污许可管理条例', '2021.03.01', '第33条：无证排污，20-100万罚款'],
    ])
    add_table_slide('山东省地方性法规', ['法规名称', '施行日期', '核心罚则'], [
        ['山东省大气污染防治条例', '2016.01.01', 'VOCs无组织排放：5万以下罚款'],
        ['山东省水污染防治条例', '2018.12.01', '超过排污许可证排放：10万-100万'],
        ['山东省环境保护条例', '2020.01.01', '按日连续处罚细化规定'],
        ['山东省固体废物污染环境防治条例', '2022.01.01', '危废规范化管理要求'],
    ])

def _build_risk_prevention_slides():
    add_section_slide('执法风险与防范')
    add_content_slide('程序合法性风险', [
        '① 超期办案：一般程序应在90日内办结，需经批准方可延长至180日',
        '② 告知不充分：处罚前未告知当事人陈述申辩权，处罚决定可能被撤销',
        '③ 听证权利未保障：当事人申请听证未依法组织的，程序违法',
        '④ 送达不规范：公告送达未满60日即申请强制执行',
    ], '建议：建立案件办理时限台账，每个节点设置提醒')
    add_content_slide('调查取证风险', [
        '① 采样不规范：无见证人、无资质、样品保存运输不符合要求',
        '② 笔录粗糙：关键信息（时间/地点/违法设备）记录不全',
        '③ 书证收集不完整：缺少环评、验收、运行台账等关键证据',
        '④ 证据链断裂：证据之间无法相互印证，孤证不能定案',
    ], '建议：制定标准化现场检查清单，逐项核查')

def _build_case_process_slides():
    add_section_slide('案件办理标准流程')
    add_content_slide('行政处罚一般程序流程图', [
        '📋 步骤1【立案】- 发现违法行为 → 填写立案审批表 → 负责人审批（24小时内）',
        '🔍 步骤2【调查】- 现场检查 + 询问调查 + 收集证据 → 形成调查报告',
        '📢 步骤3【告知】- 拟处罚决定告知书 → 送达当事人 → 听取陈述申辩（3日）',
        '⚖️ 步骤4【听证】- 当事人申请听证（3日内）→ 7日前发听证通知 → 举行听证会',
        '📝 步骤5【决定】- 重大案件集体讨论 → 法制审核 → 负责人审批签发',
        '📬 步骤6【送达】- 处罚决定书送达（7日内）→ 当事人签收',
        '💰 步骤7【执行】- 当事人自觉履行（15日内）/ 逾期加处罚款 / 申请强制执行',
        '✅ 步骤8【结案】- 执行完毕填写结案报告 → 案卷归档（保存5年）',
    ], '时限提醒：立案后90日内结案，经批准可延长至180日')
    add_content_slide('特殊程序注意事项', [
        '【移送公安】涉嫌犯罪案件：取得关键证据后5日内移交，不作行政处罚',
        '【生态损害赔偿】符合条件案件：同步启动生态环境损害赔偿调查',
        '【信息公开】处罚决定作出7日内：录入全国信用信息共享平台和国家企业信用信息公示系统',
        '【两法衔接】环境违法犯罪案件：同时抄送检察院备案，接受法律监督',
    ])

def generate_report_ppt(case_summary, output_path=None):
    """生成案件汇报 PPT"""
    new_prs()
    add_title_slide('环境违法案件查处汇报', f'案件编号：{case_summary.get("case_id","待分配")}')

    add_content_slide('一、案件基本情况', [
        f'当事人：{case_summary.get("party","（待填写）")}',
        f'案件类型：{case_summary.get("type","（待填写）")}',
        f'发现方式：{case_summary.get("discovery","专项检查")}',
        f'发现时间：{case_summary.get("date","（待填写）")}',
        f'涉及区域：{case_summary.get("location","（待填写）")}',
    ])
    add_content_slide('二、违法事实认定', [
        f'主要违法事实：{case_summary.get("fact","（待填写）")}',
        '关键证据：现场检查笔录、监测报告、询问笔录、视听资料',
        f'检测结果：{case_summary.get("result","（待填写检测数据）")}',
        f'涉案金额：{case_summary.get("amount","（待填写）")}',
    ])
    add_content_slide('三、法律适用分析', [
        f'违反法律：{case_summary.get("law","（待填写）")}',
        '违反条款：《XX法》第XX条第X款',
        '处罚依据：《XX法》第XX条第X款',
        f'裁量档次：{case_summary.get("level","（待确定）")}（依据山东省裁量基准）',
    ])
    add_content_slide('四、处理建议', [
        f'行政处罚：罚款人民币 {case_summary.get("fine","（待确定）")} 万元',
        '其他措施：责令改正 / 停产整治 / 查封扣押（根据案件情况勾选）',
        '刑事移送：{case_summary.get("criminal","否（未达追诉标准）")}',
        '下一步：提交法制审核 → 集体讨论 → 负责人审批 → 送达执行',
    ])
    add_section_slide('谢谢！')

    if not output_path:
        output_path = os.path.join(os.path.dirname(__file__), '..', 'outputs', f'案件汇报_{date.today().strftime("%Y%m%d")}.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path

if __name__ == '__main__':
    import sys
    topic = sys.argv[1] if len(sys.argv) > 1 else '行政处罚程序'
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    path = generate_training_ppt(topic, output_path)
    print(f'OK|{path}')
